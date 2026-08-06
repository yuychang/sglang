#!/usr/bin/env python3
"""Generate an illustrative deck for SGLang throughput/interactivity capacity."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "sglang-throughput-interactivity.pptx"

NAVY = RGBColor(10, 18, 32)
PANEL = RGBColor(22, 34, 53)
PANEL_2 = RGBColor(30, 46, 70)
WHITE = RGBColor(244, 248, 252)
MUTED = RGBColor(166, 184, 204)
CYAN = RGBColor(54, 211, 229)
BLUE = RGBColor(75, 130, 255)
GREEN = RGBColor(77, 218, 153)
AMBER = RGBColor(255, 184, 77)
RED = RGBColor(255, 100, 115)


def add_text(slide, x, y, w, h, text, size=20, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=18, color=WHITE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
        p.text = "•  " + p.text
    return box


def add_rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_line(slide, x1, y1, x2, y2, color=MUTED, width=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.65, 0.35, 12.0, 0.55, title, 28, WHITE, True)
    add_rect(slide, 0.65, 0.98, 1.0, 0.05, CYAN, radius=False)
    if subtitle:
        add_text(slide, 1.8, 0.88, 10.8, 0.25, subtitle, 11, MUTED)


def add_footer(slide, number):
    add_text(slide, 0.65, 7.08, 10.8, 0.18,
             "Kimi K2.5 · ISL 8K / OSL 1K · TP=4", 9, MUTED)
    add_text(slide, 12.0, 7.05, 0.65, 0.2, str(number), 9, MUTED,
             align=PP_ALIGN.RIGHT)


def new_slide(prs, title=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    if title:
        add_title(slide, title, subtitle)
    return slide


def card(slide, x, y, w, h, eyebrow, headline, body, accent=CYAN):
    add_rect(slide, x, y, w, h, PANEL)
    add_rect(slide, x, y, 0.08, h, accent, radius=False)
    add_text(slide, x + 0.25, y + 0.17, w - 0.45, 0.28,
             eyebrow.upper(), 10, accent, True)
    add_text(slide, x + 0.25, y + 0.52, w - 0.45, 0.52,
             headline, 24, WHITE, True)
    add_text(slide, x + 0.25, y + 1.08, w - 0.45, h - 1.22,
             body, 14, MUTED, valign=MSO_ANCHOR.TOP)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "SGLang Throughput Ceiling and Interactivity"
    prs.core_properties.subject = "Kimi K2.5 capacity planning on B200 and MI355X"

    # 1 — title
    slide = new_slide(prs)
    add_rect(slide, 0.7, 0.75, 0.12, 5.75, CYAN, radius=False)
    add_text(slide, 1.15, 1.05, 10.8, 1.45,
             "Calculating the Throughput Ceiling\nand Its Interactivity Cost",
             36, WHITE, True, valign=MSO_ANCHOR.TOP)
    add_text(slide, 1.15, 2.8, 10.2, 0.65,
             "An illustrative capacity model for SGLang serving",
             22, CYAN)
    add_text(slide, 1.15, 3.55, 10.2, 0.8,
             "Kimi K2.5  ·  ISL 8K / OSL 1K  ·  TP=4\n"
             "B200 + FlashInfer TRT-LLM  vs.  MI355X + AITER",
             17, MUTED, valign=MSO_ANCHOR.TOP)
    add_rect(slide, 1.15, 5.25, 10.8, 0.9, PANEL)
    add_text(slide, 1.45, 5.45, 10.2, 0.45,
             "Capacity is a frontier: maximize throughput subject to TTFT and ITL SLOs.",
             18, WHITE, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 1)

    # 2 — executive summary
    slide = new_slide(prs, "Executive summary")
    card(slide, 0.7, 1.35, 3.85, 2.1, "Memory bound", "Resident batch",
         "KV capacity determines how many long-context requests can remain active.",
         BLUE)
    card(slide, 4.75, 1.35, 3.85, 2.1, "Compute bound", "Kernel saturation",
         "Backend kernels determine when additional batching stops improving tokens/s.",
         CYAN)
    card(slide, 8.8, 1.35, 3.85, 2.1, "SLO bound", "Interactivity",
         "Larger batches increase decode-step time and therefore ITL.",
         AMBER)
    add_rect(slide, 0.7, 3.8, 11.95, 1.8, PANEL_2)
    add_text(slide, 1.05, 4.03, 11.25, 0.45,
             "The key relationship", 12, CYAN, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 4.45, 11.25, 0.55,
             "Output throughput ≈ Resident decode batch ÷ Decode-step time",
             25, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 5.0, 11.25, 0.35,
             "Throughput may rise while ITL worsens—both observations can be true.",
             15, MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3 — workload anatomy
    slide = new_slide(prs, "Workload anatomy", "Why “total throughput” can mislead")
    add_text(slide, 0.8, 1.35, 3.2, 0.5, "9,216 tokens / request", 25, WHITE, True)
    add_text(slide, 0.8, 1.95, 3.2, 0.85,
             "ISL  8,192\nOSL  1,024", 20, MUTED, True, valign=MSO_ANCHOR.TOP)
    bar_x, bar_y, bar_w, bar_h = 4.3, 1.55, 8.0, 0.72
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, PANEL_2, radius=False)
    add_rect(slide, bar_x, bar_y, bar_w * 8 / 9, bar_h, BLUE, radius=False)
    add_rect(slide, bar_x + bar_w * 8 / 9, bar_y,
             bar_w / 9, bar_h, AMBER, radius=False)
    add_text(slide, 4.55, 1.65, 5.8, 0.45, "Input / prefill  89%",
             18, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 11.42, 1.65, 0.75, 0.45, "11%",
             14, NAVY, True, align=PP_ALIGN.CENTER)
    card(slide, 0.8, 3.15, 5.65, 2.15, "Total token throughput",
         "(Input + output) / time",
         "Mostly measures prefill scaling for this 8:1 workload mix.", BLUE)
    card(slide, 6.75, 3.15, 5.55, 2.15, "Generation throughput",
         "Output tokens / time",
         "The better metric for decode capacity and ITL trade-offs.", AMBER)
    add_text(slide, 0.8, 5.78, 11.5, 0.45,
             "Report input TPS, output TPS, total TPS, TTFT, and ITL separately.",
             18, CYAN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4 — resident batch constraints
    slide = new_slide(prs, "Step 1 — Bound the resident batch")
    labels = [
        ("Offered concurrency", "Client-side request pressure", BLUE),
        ("Request slots", "Auto max_running_requests", CYAN),
        ("KV capacity", "max_total_num_tokens ÷ live tokens/request", GREEN),
        ("Resident decode batch", "Requests actually advancing on GPU", AMBER),
    ]
    widths = [11.2, 9.2, 7.2, 5.2]
    for i, ((headline, body, color), width) in enumerate(zip(labels, widths)):
        x = (13.333 - width) / 2
        y = 1.25 + i * 1.05
        add_rect(slide, x, y, width, 0.8, PANEL_2)
        add_rect(slide, x, y, 0.1, 0.8, color, radius=False)
        add_text(slide, x + 0.3, y + 0.08, width * 0.48, 0.3,
                 headline, 16, WHITE, True)
        add_text(slide, x + width * 0.5, y + 0.08, width * 0.46, 0.55,
                 body, 12, MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, 1.35, 5.68, 10.63, 0.72, PANEL)
    add_text(slide, 1.55, 5.81, 10.25, 0.42,
             "Bresident ≤ min(BKV, request slots)  ·  Offered concurrency may be much larger",
             18, WHITE, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5 — KV geometry
    slide = new_slide(prs, "Step 2 — Calculate KV capacity", "Kimi K2.5 MLA with FP8 KV")
    add_rect(slide, 0.75, 1.3, 5.75, 4.75, PANEL)
    add_text(slide, 1.05, 1.6, 5.15, 0.4, "Per-token KV geometry", 18, CYAN, True)
    add_text(slide, 1.05, 2.15, 5.15, 1.2,
             "61 layers × (512 latent + 64 RoPE)\n× 1 byte FP8",
             24, WHITE, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.25, 3.55, 4.75, 0.9, PANEL_2)
    add_text(slide, 1.45, 3.75, 4.35, 0.48,
             "= 35,136 bytes / token / GPU",
             20, GREEN, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 4.85, 5.15, 0.7,
             "TP=4 does not divide the compressed MLA KV token cost on each rank.",
             14, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.85, 1.3, 5.75, 4.75, PANEL)
    add_text(slide, 7.15, 1.6, 5.15, 0.4, "Per-request KV cost", 18, AMBER, True)
    add_text(slide, 7.15, 2.15, 5.15, 1.2,
             "35,136 bytes × (8,192 + g)",
             24, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 7.35, 3.52, 4.75, 0.4,
             "After prefill: 274.5 MiB", 18, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 7.35, 4.05, 4.75, 0.4,
             "At 1K output: 308.8 MiB", 18, AMBER, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 7.25, 4.72, 4.95, 0.75, PANEL_2)
    add_text(slide, 7.45, 4.88, 4.55, 0.4,
             "BKV,end = floor(TKV / 9,216)",
             18, WHITE, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 5)

    # 6 — platform comparison
    slide = new_slide(prs, "Step 3 — Compare raw memory headroom",
                      "Illustrative upper bounds before runtime reservations")
    add_text(slide, 0.75, 1.2, 5.8, 0.35,
             "B200 · NVIDIA Kimi-K2.5-NVFP4", 17, BLUE, True)
    add_text(slide, 6.85, 1.2, 5.7, 0.35,
             "MI355X · AMD Kimi-K2.5-MXFP4", 17, GREEN, True)
    # stacked memory bars scaled to 288
    scale = 3.8 / 288
    for x, hbm, weights, accent in [
        (1.15, 180, 137.55, BLUE),
        (7.25, 288, 130.15, GREEN),
    ]:
        total_h = hbm * scale
        weight_h = weights * scale
        add_rect(slide, x, 5.55 - total_h, 2.2, total_h, PANEL_2, radius=False)
        add_rect(slide, x, 5.55 - weight_h, 2.2, weight_h, accent, radius=False)
        add_text(slide, x, 5.62, 2.2, 0.32, f"{hbm:.0f} GiB HBM",
                 13, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.05, 5.55 - weight_h + 0.1, 2.1, 0.45,
                 f"{weights:.1f} GiB\nweights", 13, WHITE, True,
                 align=PP_ALIGN.CENTER)
        headroom = hbm - weights
        add_text(slide, x + 2.45, 5.55 - total_h + 0.1, 2.45, 0.65,
                 f"Raw headroom\n{headroom:.1f} GiB", 15, accent, True)
        add_line(slide, x + 2.2, 5.55 - total_h + 0.35,
                 x + 2.4, 5.55 - total_h + 0.35, accent, 2)
    add_rect(slide, 3.9, 2.35, 2.85, 1.55, PANEL)
    add_text(slide, 4.12, 2.53, 2.4, 0.3, "Raw 9K KV bound", 12, CYAN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 4.12, 2.95, 2.4, 0.55, "B200  ≈ 141", 21, BLUE, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 4.12, 3.42, 2.4, 0.35, "MI355X  ≈ 523", 21, GREEN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 0.85, 6.35, 11.75, 0.4,
             "Use /server_info max_total_num_tokens for the real bound; allocator, graph, VLM, and backend reserves reduce these raw values.",
             12, MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7 — throughput vs ITL
    slide = new_slide(prs, "Step 4 — Translate batch size into throughput and ITL")
    add_rect(slide, 0.75, 1.25, 4.05, 4.95, PANEL)
    add_text(slide, 1.05, 1.62, 3.45, 0.55,
             "Output TPS(B) ≈ B / t(B)", 24, CYAN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 2.35, 3.45, 0.55,
             "ITL(B) ≈ t(B)", 24, AMBER, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 3.25, 3.45, 1.6,
             "If batch grows faster than step time:\n\n"
             "• throughput increases\n"
             "• ITL also increases",
             18, WHITE, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.15, 1.25, 7.45, 4.95, PANEL)
    add_text(slide, 5.5, 1.55, 6.75, 0.4,
             "Illustrative normalized example", 15, WHITE, True,
             align=PP_ALIGN.CENTER)
    # axes
    add_line(slide, 6.0, 5.35, 11.9, 5.35, MUTED, 1)
    add_line(slide, 6.0, 2.15, 6.0, 5.35, MUTED, 1)
    points = [(6.4, 4.55, "128"), (8.0, 3.78, "256"),
              (9.65, 3.15, "512"), (11.35, 2.75, "1024")]
    for idx, (x, y, label) in enumerate(points):
        if idx:
            px, py, _ = points[idx - 1]
            add_line(slide, px, py, x, y, CYAN, 3)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.08),
                                     Inches(y - 0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = CYAN
        dot.line.color.rgb = CYAN
        add_text(slide, x - 0.35, 5.45, 0.7, 0.25, label, 10, MUTED,
                 align=PP_ALIGN.CENTER)
    add_text(slide, 5.15, 5.82, 7.45, 0.25,
             "Offered concurrency →", 11, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, 5.35, 2.0, 0.5, 2.8, "TPS\n↑", 11, CYAN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 7.0, 2.3, 4.8, 0.55,
             "More throughput\nbut longer decode steps", 17, AMBER, True,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8 — why curves differ
    slide = new_slide(prs, "Why B200 shows a ceiling while MI355X keeps scaling")
    card(slide, 0.7, 1.25, 3.85, 2.15, "1 · Memory",
         "~3.7× raw KV headroom",
         "MI355X can sustain a much larger resident long-context batch.", GREEN)
    card(slide, 4.75, 1.25, 3.85, 2.15, "2 · Kernels",
         "Different saturation points",
         "FlashInfer TRT-LLM reaches useful occupancy earlier; AITER persistent MLA can benefit from larger batches.",
         CYAN)
    card(slide, 8.8, 1.25, 3.85, 2.15, "3 · Metric mix",
         "Prefill dominates total TPS",
         "At 8K/1K, improved AITER prefill batching can raise total TPS while ITL deteriorates.",
         BLUE)
    add_rect(slide, 0.7, 3.78, 11.95, 1.95, PANEL_2)
    add_text(slide, 1.0, 4.0, 2.4, 0.35, "B200 pattern", 13, BLUE, True)
    add_text(slide, 3.1, 3.95, 8.9, 0.45,
             "Capacity / compute knee near ~128 → total TPS flattens", 18, WHITE, True)
    add_text(slide, 1.0, 4.75, 2.4, 0.35, "MI355X pattern", 13, GREEN, True)
    add_text(slide, 3.1, 4.7, 8.9, 0.72,
             "Resident batch scales into the hundreds; at 1024 offered concurrency, queueing and prefill batching may still raise total TPS.",
             18, WHITE, True)
    add_footer(slide, 8)

    # 9 — end-to-end upper bound
    slide = new_slide(prs, "Step 5 — Bound end-to-end throughput")
    add_rect(slide, 0.8, 1.25, 5.7, 2.15, PANEL)
    add_text(slide, 1.05, 1.52, 5.2, 0.35,
             "No prefill/decode overlap", 14, CYAN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 2.0, 5.2, 0.72,
             "QPS ≤ 1 / (ISL / P + OSL / D)",
             23, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 2.77, 5.2, 0.35,
             "Total TPS ≤ 9,216 × QPS", 16, MUTED,
             align=PP_ALIGN.CENTER)
    add_rect(slide, 6.85, 1.25, 5.7, 2.15, PANEL)
    add_text(slide, 7.1, 1.52, 5.2, 0.35,
             "Ideal independent overlap", 14, GREEN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 7.1, 2.0, 5.2, 0.72,
             "QPS ≤ min(P / ISL, D / OSL)",
             23, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 7.1, 2.77, 5.2, 0.35,
             "Real mixed scheduling lies between the bounds", 16, MUTED,
             align=PP_ALIGN.CENTER)
    # pipeline diagram
    add_rect(slide, 1.0, 4.2, 2.55, 0.8, BLUE)
    add_text(slide, 1.15, 4.33, 2.25, 0.45, "Queue", 18, WHITE, True,
             align=PP_ALIGN.CENTER)
    add_line(slide, 3.55, 4.6, 4.25, 4.6, MUTED, 3)
    add_rect(slide, 4.25, 4.2, 2.55, 0.8, CYAN)
    add_text(slide, 4.4, 4.33, 2.25, 0.45, "Prefill · P", 18, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_line(slide, 6.8, 4.6, 7.5, 4.6, MUTED, 3)
    add_rect(slide, 7.5, 4.2, 2.55, 0.8, AMBER)
    add_text(slide, 7.65, 4.33, 2.25, 0.45, "Decode · D", 18, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_line(slide, 10.05, 4.6, 10.75, 4.6, MUTED, 3)
    add_rect(slide, 10.75, 4.2, 1.55, 0.8, GREEN)
    add_text(slide, 10.9, 4.33, 1.25, 0.45, "Done", 18, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 1.0, 5.45, 11.3, 0.65,
             "P and D should be measured at steady state; a hardware-only roofline is rarely predictive for distributed MoE serving.",
             15, MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10 — telemetry
    slide = new_slide(prs, "Measure the bound instead of guessing it")
    add_rect(slide, 0.75, 1.25, 7.05, 4.95, PANEL)
    add_text(slide, 1.05, 1.55, 6.45, 0.35,
             "Live SGLang signals", 17, CYAN, True)
    metrics = [
        ("num_running_reqs", "Actual resident scheduler batch"),
        ("num_waiting_reqs", "Offered load beyond admission"),
        ("token_usage", "KV pressure"),
        ("gen_throughput", "Decode—not total—throughput"),
        ("cache_hit_rate", "Shared-prefix correction"),
        ("queues.retracted", "KV churn and likely ITL spikes"),
    ]
    for i, (name, meaning) in enumerate(metrics):
        y = 2.02 + i * 0.62
        add_text(slide, 1.05, y, 2.65, 0.3, name, 13, WHITE, True,
                 font="Aptos Mono")
        add_text(slide, 3.78, y, 3.65, 0.3, meaning, 13, MUTED)
    add_rect(slide, 8.15, 1.25, 4.45, 4.95, PANEL_2)
    add_text(slide, 8.45, 1.58, 3.85, 0.35,
             "Query during the sweep", 16, GREEN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 8.47, 2.15, 3.8, 2.4,
             "curl -s \\\n"
             "  'http://host:30000/\n"
             "   v1/loads?\n"
             "   include=core,queues'\n"
             "  | jq",
             15, WHITE, False, font="Aptos Mono",
             valign=MSO_ANCHOR.TOP)
    add_text(slide, 8.45, 5.15, 3.85, 0.55,
             "Capture p50 and p99—not only averages.",
             15, AMBER, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    # 11 — decision frontier
    slide = new_slide(prs, "Choose the SLO knee—not maximum concurrency")
    add_rect(slide, 0.75, 1.25, 7.1, 4.95, PANEL)
    add_line(slide, 1.45, 5.35, 7.15, 5.35, MUTED, 1)
    add_line(slide, 1.45, 1.95, 1.45, 5.35, MUTED, 1)
    curve = [(1.9, 4.65), (2.8, 3.75), (3.8, 3.0),
             (4.85, 2.55), (5.9, 2.4), (6.75, 2.37)]
    for i, (x, y) in enumerate(curve):
        if i:
            px, py = curve[i - 1]
            add_line(slide, px, py, x, y, CYAN, 4)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.07),
                                     Inches(y - 0.07), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = CYAN
        dot.line.color.rgb = CYAN
    knee_x, knee_y = curve[3]
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(knee_x - 0.18),
                                  Inches(knee_y - 0.18), Inches(0.36), Inches(0.36))
    ring.fill.background()
    ring.line.color.rgb = AMBER
    ring.line.width = Pt(3)
    add_text(slide, 5.05, 2.72, 1.85, 0.5, "SLO knee", 15, AMBER, True)
    add_text(slide, 1.5, 5.55, 5.5, 0.28, "Concurrency / ITL cost →", 11, MUTED,
             align=PP_ALIGN.CENTER)
    add_text(slide, 0.82, 2.0, 0.5, 2.8, "TPS\n↑", 11, CYAN, True,
             align=PP_ALIGN.CENTER)
    add_rect(slide, 8.2, 1.25, 4.4, 4.95, PANEL_2)
    add_text(slide, 8.55, 1.62, 3.7, 0.45,
             "Optimization objective", 17, WHITE, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 8.55, 2.35, 3.7, 1.1,
             "Maximize output TPS",
             25, GREEN, True, align=PP_ALIGN.CENTER)
    add_text(slide, 8.55, 3.35, 3.7, 0.4,
             "subject to", 14, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, 8.55, 3.92, 3.7, 1.1,
             "p99 ITL ≤ SLO\np99 TTFT ≤ SLO",
             21, AMBER, True, align=PP_ALIGN.CENTER)
    add_text(slide, 8.55, 5.28, 3.7, 0.45,
             "Capacity is workload-specific.", 14, MUTED,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12 — takeaways
    slide = new_slide(prs, "Takeaways")
    takeaways = [
        ("1", "Separate the limits",
         "Request slots, KV capacity, kernel saturation, and SLOs are different constraints.",
         BLUE),
        ("2", "Use the right metric",
         "At ISL 8K / OSL 1K, total TPS mostly measures prefill; use output TPS for decode.",
         CYAN),
        ("3", "Explain the platform gap",
         "MI355X has far more KV headroom and AITER scales to larger batches; B200 reaches its knee earlier.",
         GREEN),
        ("4", "Optimize for interactivity",
         "Choose the highest output throughput that satisfies p99 TTFT and ITL.",
         AMBER),
    ]
    for i, (num, headline, body, accent) in enumerate(takeaways):
        y = 1.25 + i * 1.35
        add_rect(slide, 0.85, y, 11.65, 1.05, PANEL)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1),
                                        Inches(y + 0.17), Inches(0.7), Inches(0.7))
        circle.fill.solid(); circle.fill.fore_color.rgb = accent
        circle.line.color.rgb = accent
        add_text(slide, 1.1, y + 0.17, 0.7, 0.7, num, 18, NAVY, True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, 2.1, y + 0.12, 3.1, 0.35, headline, 18, WHITE, True)
        add_text(slide, 5.1, y + 0.12, 6.95, 0.72, body, 14, MUTED)
    add_footer(slide, 12)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
